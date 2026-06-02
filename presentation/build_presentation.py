"""Generate LeakLens PFA PowerPoint presentation."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIGURES = ROOT / "rapport" / "figures"
OUT = Path(__file__).resolve().parent / "LeakLens-PFA-Presentation.pptx"

# Theme (matches iTeam / cyber purple template)
PURPLE = RGBColor(0x5E, 0x50, 0xA1)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF4, 0xF8)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=20, color=DARK, spacing=Pt(8)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
    return box


def add_footer_bar(slide, title: str):
    add_rect(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), PURPLE)
    add_textbox(slide, Inches(0.4), Inches(7.08), Inches(8), Inches(0.35), title, size=11, bold=True, color=WHITE)
    add_textbox(slide, Inches(10.5), Inches(7.08), Inches(2.5), Inches(0.35), "LeakLens PFA 2026", size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def slide_header(slide, title: str, subtitle: str = ""):
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, PURPLE)
    add_textbox(slide, Inches(0.45), Inches(0.35), Inches(12), Inches(0.6), title, size=32, bold=True, color=PURPLE)
    if subtitle:
        add_textbox(slide, Inches(0.45), Inches(0.95), Inches(12), Inches(0.4), subtitle, size=14, color=GRAY)
    add_rect(slide, Inches(0.45), Inches(1.35), Inches(12.4), Inches(0.03), PURPLE)


def add_image_safe(slide, path: Path, left, top, width, height=None):
    if path.exists():
        if height:
            slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(str(path), left, top, width=width)
        return True
    add_textbox(slide, left, top, width, Inches(0.5), f"[Missing: {path.name}]", size=12, color=GRAY)
    return False


def build_title_slide(prs: Presentation):
    slide = blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_BG)

    logo = FIGURES / "iteam-logo.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.5), Inches(0.35), height=Inches(0.85))

    add_textbox(
        slide, Inches(3.2), Inches(0.45), Inches(9.5), Inches(0.55),
        "République Tunisienne — Ministère de l'Enseignement Supérieur\net de la Recherche Scientifique",
        size=11, color=GRAY, align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.5),
        "END-OF-YEAR PROJECT PRESENTATION",
        size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(0.8), Inches(2.75), Inches(11.7), Inches(1.4),
        "OSINT INTELLIGENCE PLATFORM\nFOR LEAKED DATA RESEARCH",
        size=36, bold=True, color=PURPLE, align=PP_ALIGN.CENTER,
    )

    add_textbox(slide, Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.45),
                "LeakLens — Cybersecurity · iTeam University", size=16, color=GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(5.35), Inches(5), Inches(0.35),
                "Prepared by:", size=13, color=GRAY)
    add_textbox(slide, Inches(0.8), Inches(5.65), Inches(5), Inches(0.45),
                "Motez Machghoul", size=18, bold=True, color=DARK)

    add_textbox(slide, Inches(7.5), Inches(5.35), Inches(5), Inches(0.35),
                "Supervised by:", size=13, color=GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(7.5), Inches(5.65), Inches(5), Inches(0.45),
                "Mr. Raihane Modhaffer", size=18, bold=True, color=DARK, align=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(0), Inches(6.55), Inches(10.5), Inches(0.55), PURPLE)
    add_textbox(slide, Inches(0.5), Inches(6.62), Inches(6), Inches(0.4),
                "Academic Year: 2025 – 2026", size=14, bold=True, color=WHITE)

    # Decorative accent (right)
    add_rect(slide, Inches(12.6), Inches(5.8), Inches(0.35), Inches(1.7), PURPLE)
    add_rect(slide, Inches(12.95), Inches(6.1), Inches(0.35), Inches(1.4), ACCENT)


def build_plan_slide(prs: Presentation):
    slide = blank_slide(prs)
    slide_header(slide, "Presentation Plan")

    plan = [
        ("01", "Introduction", "Why breach data matters in cybersecurity"),
        ("02", "Problem Statement", "The exposure gap defenders face"),
        ("03", "Proposed Solution", "LeakLens — our OSINT platform"),
        ("04", "System Design", "Architecture, search flow, security"),
        ("05", "Implementation", "UI, Docker deployment, live demos"),
        ("06", "Results & Conclusion", "What we achieved — Q&A"),
    ]

    y = Inches(1.75)
    for num, title, desc in plan:
        add_rect(slide, Inches(0.55), y, Inches(0.65), Inches(0.65), PURPLE)
        add_textbox(slide, Inches(0.55), y + Inches(0.12), Inches(0.65), Inches(0.4),
                    num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.4), y + Inches(0.02), Inches(4), Inches(0.35),
                    title, size=20, bold=True, color=DARK)
        add_textbox(slide, Inches(1.4), y + Inches(0.38), Inches(10), Inches(0.35),
                    desc, size=14, color=GRAY)
        y += Inches(0.85)

    add_footer_bar(slide, "LeakLens — End-of-Year Project Defense")


def build_content_slide(prs, title, subtitle, bullets, footer):
    slide = blank_slide(prs)
    slide_header(slide, title, subtitle)
    add_bullets(slide, Inches(0.55), Inches(1.65), Inches(12.2), Inches(5), bullets, size=22)
    add_footer_bar(slide, footer)


def build_image_slide(prs, title, subtitle, image_name, caption, footer, img_width=Inches(10.5)):
    slide = blank_slide(prs)
    slide_header(slide, title, subtitle)
    left = (SLIDE_W - img_width) / 2
    add_image_safe(slide, FIGURES / image_name, left, Inches(1.55), img_width, height=Inches(4.6))
    if caption:
        add_textbox(slide, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.5), caption, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer_bar(slide, footer)


def build_two_image_slide(prs, title, subtitle, img1, img2, caption, footer):
    slide = blank_slide(prs)
    slide_header(slide, title, subtitle)
    add_image_safe(slide, FIGURES / img1, Inches(0.45), Inches(1.55), Inches(6.1), height=Inches(4.2))
    add_image_safe(slide, FIGURES / img2, Inches(6.75), Inches(1.55), Inches(6.1), height=Inches(4.2))
    add_textbox(slide, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.5), caption, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer_bar(slide, footer)


def build_demo_slide(prs, title, stats, image_name, footer):
    slide = blank_slide(prs)
    slide_header(slide, title, "Live query on production — leakslens.vercel.app")
    add_bullets(slide, Inches(0.55), Inches(1.6), Inches(5.5), Inches(4.5), stats, size=18)
    add_image_safe(slide, FIGURES / image_name, Inches(6.3), Inches(1.55), Inches(6.5), height=Inches(4.5))
    add_footer_bar(slide, footer)


def build_thank_you(prs):
    slide = blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PURPLE)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1),
                "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.6),
                "Questions & Discussion", size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
                "leakslens.vercel.app  ·  Motez Machghoul  ·  iTeam University 2026",
                size=14, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_plan_slide(prs)

    build_content_slide(
        prs,
        "Introduction",
        "Open-Source Intelligence (OSINT) for breach exposure assessment",
        [
            "Every day, billions of credentials leak onto forums, stealer logs, and combolists.",
            "Attackers use this data for credential stuffing, account takeover, and phishing.",
            "Defenders need the same visibility — to find exposure before attackers exploit it.",
            "LeakLens was built to give authorized analysts a secure way to search breach data.",
        ],
        "01 — Introduction",
    )

    build_content_slide(
        prs,
        "Problem Statement",
        "The defender's blind spot",
        [
            "Breach data stays online long after an incident is announced.",
            "Reused passwords turn one small leak into access to email, cloud, and VPN accounts.",
            "Commercial OSINT tools are expensive, closed, or limited to breach names only.",
            "Students and small security teams lack an affordable, transparent research platform.",
        ],
        "02 — Problem Statement",
    )

    build_content_slide(
        prs,
        "Proposed Solution — LeakLens",
        "An OSINT breach-search platform built for cybersecurity research",
        [
            "Search 7 types: email, username, IP, password, name, hash, domain.",
            "Return real breach records with severity scoring and enrichment context.",
            "Protect access: invite-only signup, JWT auth, RBAC, rate limiting.",
            "Deploy transparently: Next.js on Vercel + FastAPI on Docker VPS.",
        ],
        "03 — Proposed Solution",
    )

    build_image_slide(
        prs,
        "Who Uses the Platform",
        "Role-based access control (RBAC)",
        "usecase-diagram.png",
        "Guest registers with invite code · Analyst searches · Admin manages users",
        "04 — System Design",
    )

    build_image_slide(
        prs,
        "Global Architecture",
        "Split deployment — frontend on Vercel, backend on VPS",
        "architecture-diag.png",
        "Browser → Next.js → Nginx Proxy Manager → Docker (FastAPI + PostgreSQL) → Snusbase & enrichment APIs",
        "04 — System Design",
    )

    build_image_slide(
        prs,
        "Search Engine Pipeline",
        "From user input to breach results in seconds",
        "search-sequence.png",
        "POST /api/v1/search → validate → Search Engine → Snusbase + API Services → PostgreSQL audit → response",
        "04 — System Design",
    )

    build_content_slide(
        prs,
        "Security by Design",
        "Defense-in-depth for a sensitive OSINT capability",
        [
            "JWT access (15 min) + refresh tokens · bcrypt password hashing.",
            "Roles: admin, analyst, viewer — search restricted to analyst/admin.",
            "Rate limits on auth and search · account lockout after failed logins.",
            "Security headers: CSP, HSTS, CORS · no local storage of raw credentials.",
        ],
        "04 — Security",
    )

    build_two_image_slide(
        prs,
        "User Interface",
        "Dark-themed UI for security analysts",
        "landing-guest.png",
        "home-auth.png",
        "Public landing page and authenticated search workspace with 7 search types",
        "05 — Implementation",
    )

    build_two_image_slide(
        prs,
        "Docker Deployment on VPS",
        "Containerized backend — Alpine, non-root, health checks",
        "docker-build.png",
        "docker-ps.png",
        "docker compose up -d --build · backend + PostgreSQL containers healthy on port 8500",
        "05 — Implementation",
    )

    build_demo_slide(
        prs,
        "Demo — Password Search",
        [
            "Query: qwerty@",
            "Type: Password search",
            "Result: 748 leaked records",
            "Exposure: emails, usernames, plaintext passwords",
            "Severity: mostly CRITICAL",
            "Insight: weak passwords spread across many breaches",
        ],
        "password-search.png",
        "06 — Live Demonstration",
    )

    build_demo_slide(
        prs,
        "Demo — Domain Search",
        [
            "Query: itbs.tn",
            "Type: Domain search",
            "Result: 16 breached student records",
            "Context: institutional email exposure",
            "Enrichment: subdomains + WHOIS data",
            "Use case: campus security & incident response",
        ],
        "domain-search.png",
        "06 — Live Demonstration",
    )

    build_demo_slide(
        prs,
        "Demo — Name Search",
        [
            "Query: personal full name",
            "Type: Name search",
            "Result: matching leak records found",
            "Exposure: linked email + plaintext password",
            "Use case: individual exposure check",
            "Anyone can verify if their identity appears in leaks",
        ],
        "name-search.png",
        "06 — Live Demonstration",
    )

    build_content_slide(
        prs,
        "What We Achieved",
        "From problem to production-ready platform",
        [
            "Built a full-stack OSINT platform deployed at leakslens.vercel.app.",
            "Integrated Snusbase + enrichment APIs with concurrent search engine.",
            "Implemented JWT, RBAC, invite gate, and Docker-hardened backend.",
            "Validated with real queries: 748 password hits, 16 domain records, personal name exposure.",
        ],
        "07 — Conclusion",
    )

    build_thank_you(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
